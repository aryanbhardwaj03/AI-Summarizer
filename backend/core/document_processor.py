"""
Universal Document Processor.
Handles PDF, DOCX, PPTX, TXT, and Markdown files.
"""

import io
import logging
from pathlib import Path

from langchain_text_splitters import RecursiveCharacterTextSplitter

from config import settings

logger = logging.getLogger(__name__)


class DocumentChunk:
    """A chunk of text with metadata."""

    def __init__(self, text: str, metadata: dict):
        self.text = text
        self.metadata = metadata


class UniversalDocumentProcessor:
    """Processes any supported document type into text chunks."""

    def __init__(self):
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=settings.CHUNK_SIZE,
            chunk_overlap=settings.CHUNK_OVERLAP,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""],
        )

    def process(self, file_path: str, file_type: str) -> dict:
        """
        Process any supported document type.

        Returns:
            dict with 'chunks', 'images', 'tables', 'metadata'
        """
        file_path = Path(file_path)
        logger.info(f"Processing {file_type.upper()}: {file_path.name}")

        processors = {
            "pdf": self._process_pdf,
            "docx": self._process_docx,
            "pptx": self._process_pptx,
            "txt": self._process_txt,
            "md": self._process_markdown,
        }

        processor = processors.get(file_type)
        if not processor:
            raise ValueError(f"Unsupported file type: {file_type}")

        return processor(file_path)

    def _process_pdf(self, file_path: Path) -> dict:
        """Process a PDF file."""
        import fitz
        import pdfplumber
        from PIL import Image

        result = {"chunks": [], "images": [], "tables": [], "metadata": {"total_pages": 0}}

        # Extract text with PyMuPDF
        doc = fitz.open(str(file_path))
        result["metadata"]["total_pages"] = len(doc)
        all_text = {}

        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text")
            all_text[page_num + 1] = text

            # Extract images
            for img_idx, img_info in enumerate(page.get_images(full=True)):
                try:
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    image = Image.open(io.BytesIO(base_image["image"]))
                    if image.width >= 50 and image.height >= 50:
                        result["images"].append({
                            "image": image,
                            "page": page_num + 1,
                            "index": img_idx,
                            "size": (image.width, image.height),
                        })
                except Exception:
                    pass

        doc.close()

        # Skipping pdfplumber table extraction to massively speed up processing on large documents
        # tables will be handled via raw text in fitz

        # Chunk text
        for page_num, text in all_text.items():
            if not text.strip():
                continue
            for i, chunk_text in enumerate(self.splitter.split_text(text)):
                result["chunks"].append(DocumentChunk(
                    text=chunk_text,
                    metadata={"page": page_num, "chunk_index": i, "type": "text"},
                ))

        # Add tables as chunks
        for table in result["tables"]:
            result["chunks"].append(DocumentChunk(
                text=table["text"],
                metadata={"page": table["page"], "chunk_index": 0, "type": "table"},
            ))

        return result

    def _process_docx(self, file_path: Path) -> dict:
        """Process a DOCX file."""
        from docx import Document as DocxDocument

        doc = DocxDocument(str(file_path))
        result = {"chunks": [], "images": [], "tables": [], "metadata": {"total_pages": 0}}

        # Extract text from paragraphs
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)

        # Extract tables
        for table in doc.tables:
            rows = []
            for row in table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            if rows:
                md = self._table_to_markdown(rows)
                result["tables"].append({"text": md, "page": 0})

        # Estimate pages
        combined = "\n".join(full_text)
        result["metadata"]["total_pages"] = max(1, len(combined) // 3000)

        # Chunk
        for i, chunk_text in enumerate(self.splitter.split_text(combined)):
            result["chunks"].append(DocumentChunk(
                text=chunk_text,
                metadata={"page": 0, "chunk_index": i, "type": "text"},
            ))

        for table in result["tables"]:
            result["chunks"].append(DocumentChunk(
                text=table["text"],
                metadata={"page": 0, "chunk_index": 0, "type": "table"},
            ))

        return result

    def _process_pptx(self, file_path: Path) -> dict:
        """Process a PowerPoint file."""
        from pptx import Presentation

        prs = Presentation(str(file_path))
        result = {"chunks": [], "images": [], "tables": [], "metadata": {"total_pages": len(prs.slides)}}

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_text.append(para.text)

                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        rows.append([cell.text.strip() for cell in row.cells])
                    if rows:
                        md = self._table_to_markdown(rows)
                        result["tables"].append({"text": md, "page": slide_num})

            if slide_text:
                text = f"Slide {slide_num}:\n" + "\n".join(slide_text)
                for i, chunk_text in enumerate(self.splitter.split_text(text)):
                    result["chunks"].append(DocumentChunk(
                        text=chunk_text,
                        metadata={"page": slide_num, "chunk_index": i, "type": "slide"},
                    ))

        for table in result["tables"]:
            result["chunks"].append(DocumentChunk(
                text=table["text"],
                metadata={"page": table["page"], "chunk_index": 0, "type": "table"},
            ))

        return result

    def _process_txt(self, file_path: Path) -> dict:
        """Process a plain text file."""
        text = file_path.read_text(encoding="utf-8", errors="ignore")
        result = {"chunks": [], "images": [], "tables": [], "metadata": {"total_pages": max(1, len(text) // 3000)}}

        for i, chunk_text in enumerate(self.splitter.split_text(text)):
            result["chunks"].append(DocumentChunk(
                text=chunk_text,
                metadata={"page": 0, "chunk_index": i, "type": "text"},
            ))

        return result

    def _process_markdown(self, file_path: Path) -> dict:
        """Process a Markdown file."""
        import markdown
        text = file_path.read_text(encoding="utf-8", errors="ignore")

        # Strip markdown to plain text for chunking
        # Keep the raw markdown for display
        result = {"chunks": [], "images": [], "tables": [], "metadata": {"total_pages": max(1, len(text) // 3000)}}

        for i, chunk_text in enumerate(self.splitter.split_text(text)):
            result["chunks"].append(DocumentChunk(
                text=chunk_text,
                metadata={"page": 0, "chunk_index": i, "type": "markdown"},
            ))

        return result

    @staticmethod
    def _table_to_markdown(table: list) -> str:
        """Convert table to markdown."""
        if not table or not table[0]:
            return ""

        def clean(cell):
            return str(cell).strip().replace("|", "\\|").replace("\n", " ") if cell else ""

        header = [clean(c) for c in table[0]]
        num_cols = len(header)

        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(["---"] * num_cols) + " |",
        ]
        for row in table[1:]:
            cleaned = [clean(c) for c in row]
            while len(cleaned) < num_cols:
                cleaned.append("")
            lines.append("| " + " | ".join(cleaned[:num_cols]) + " |")

        return "\n".join(lines)
